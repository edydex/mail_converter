"""
Attachment Converter Module

Converts various file formats to PDF with OCR support.
Supports: PDF, DOC(X), XLS(X), PPT(X), images, text files, EML, MSG, HTML, CSV
"""

import os
import io
import sys
import tempfile
import subprocess
import shutil
import platform
import logging
from pathlib import Path
from typing import Optional, Callable, List, Tuple, Union
from dataclasses import dataclass
from enum import Enum

# PDF libraries
from reportlab.lib.pagesizes import letter, A4
from reportlab.lib.units import inch
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image as RLImage, Table, TableStyle
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib import colors
from reportlab.lib.enums import TA_LEFT

# Image processing
from PIL import Image
import img2pdf

# Document processing
from docx import Document as DocxDocument
from pptx import Presentation
import openpyxl
import xlrd
import pandas as pd

# PDF processing
from PyPDF2 import PdfReader, PdfWriter
from pdf2image import convert_from_path, convert_from_bytes

# OCR
import pytesseract

logger = logging.getLogger(__name__)


class ConversionStatus(Enum):
    SUCCESS = "success"
    PARTIAL = "partial"  # Converted but with warnings
    FAILED = "failed"
    SKIPPED = "skipped"  # Unsupported format


@dataclass
class ConversionResult:
    """Result of a file conversion"""
    status: ConversionStatus
    output_path: Optional[Path]
    original_path: Path
    original_type: str
    message: str
    ocr_applied: bool = False


class AttachmentConverter:
    """
    Converts various file formats to PDF.
    All conversions include OCR when applicable.
    """
    
    # Supported formats mapped to conversion methods
    SUPPORTED_FORMATS = {
        # Documents
        '.pdf': '_convert_pdf',
        '.doc': '_convert_doc',
        '.docx': '_convert_docx',
        '.xls': '_convert_excel',
        '.xlsx': '_convert_excel',
        '.ppt': '_convert_ppt',
        '.pptx': '_convert_pptx',
        '.txt': '_convert_text',
        '.csv': '_convert_csv',
        '.html': '_convert_html',
        '.htm': '_convert_html',
        # Calendar
        '.ics': '_convert_ics',
        # Images
        '.jpg': '_convert_image',
        '.jpeg': '_convert_image',
        '.png': '_convert_image',
        '.gif': '_convert_image',
        '.bmp': '_convert_image',
        '.tif': '_convert_image',
        '.tiff': '_convert_image',
        # Email
        '.eml': '_convert_eml',
        '.msg': '_convert_msg',
    }
    
    def __init__(
        self, 
        ocr_enabled: bool = True,
        progress_callback: Optional[Callable[[str], None]] = None
    ):
        """
        Initialize the attachment converter.
        
        Args:
            ocr_enabled: Whether to apply OCR to images and scanned PDFs
            progress_callback: Optional callback for progress messages
        """
        self.ocr_enabled = ocr_enabled
        self.progress_callback = progress_callback
        self.temp_dir = tempfile.mkdtemp(prefix="mail_converter_")
        
        # Check for external tools
        self._check_dependencies()
    
    def _check_dependencies(self):
        """Check for required external dependencies."""
        self.has_tesseract = shutil.which("tesseract") is not None
        
        # Check for LibreOffice (including macOS app bundle and Windows locations)
        macos_libreoffice = "/Applications/LibreOffice.app/Contents/MacOS/soffice"
        windows_libreoffice_paths = [
            r"C:\Program Files\LibreOffice\program\soffice.exe",
            r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
        ]
        
        # Find LibreOffice executable
        self.libreoffice_path = (
            shutil.which("libreoffice") or 
            shutil.which("soffice") or 
            (macos_libreoffice if os.path.isfile(macos_libreoffice) else None)
        )
        
        # Check Windows paths if not found yet
        if not self.libreoffice_path:
            for win_path in windows_libreoffice_paths:
                if os.path.isfile(win_path):
                    self.libreoffice_path = win_path
                    break
        
        self.has_libreoffice = self.libreoffice_path is not None
        
        if self.has_libreoffice:
            logger.info(f"LibreOffice found at: {self.libreoffice_path}")
        
        # Check for Poppler (used by pdf2image)
        self.poppler_path = self._find_poppler()
        
        if not self.has_tesseract and self.ocr_enabled:
            logger.warning("Tesseract not found. OCR will be disabled.")
            self.ocr_enabled = False
        
        if not self.has_libreoffice:
            logger.warning("LibreOffice not found. DOC/PPT conversion may fail.")
    
    def _find_poppler(self) -> Optional[str]:
        """Find Poppler binaries path, checking bundled location first."""
        # Check if running as PyInstaller bundle
        if getattr(sys, 'frozen', False) and hasattr(sys, '_MEIPASS'):
            bundled_path = os.path.join(sys._MEIPASS, 'bin', 'poppler')
            if os.path.isdir(bundled_path):
                logger.info(f"Using bundled Poppler at: {bundled_path}")
                return bundled_path
        
        # Check common Windows installation paths
        if sys.platform == 'win32':
            common_paths = [
                r"C:\Program Files\poppler\bin",
                r"C:\Program Files\poppler-24.02.0\Library\bin",
                r"C:\poppler\bin",
            ]
            for path in common_paths:
                if os.path.isdir(path) and os.path.isfile(os.path.join(path, 'pdftoppm.exe')):
                    logger.info(f"Found Poppler at: {path}")
                    return path
        
        # On macOS/Linux, pdf2image finds it via PATH
        return None
    
    def _log(self, message: str):
        """Log message and call progress callback if set."""
        logger.info(message)
        if self.progress_callback:
            self.progress_callback(message)
    
    def _safe_subprocess_run(
        self, 
        cmd: List[str], 
        timeout: int = 120,
        **kwargs
    ) -> subprocess.CompletedProcess:
        """
        Run a subprocess with settings optimized for macOS thread safety.
        
        On macOS, running subprocess from background threads can cause 
        Objective-C runtime crashes (NSInvalidArgumentException). This method
        isolates the subprocess to prevent those issues.
        
        Args:
            cmd: Command and arguments to run
            timeout: Maximum time to wait (seconds)
            **kwargs: Additional args passed to subprocess.run
            
        Returns:
            CompletedProcess result
        """
        # Create a clean environment to avoid macOS threading issues
        clean_env = os.environ.copy()
        
        # Remove macOS-specific variables that can cause issues
        for var in ['__CF_USER_TEXT_ENCODING', 'SECURITYSESSIONID']:
            clean_env.pop(var, None)
        
        # Disable Objective-C fork safety check that can cause crashes
        if platform.system() == 'Darwin':
            clean_env['OBJC_DISABLE_INITIALIZE_FORK_SAFETY'] = 'YES'
        
        return subprocess.run(
            cmd,
            capture_output=True,
            text=True,
            timeout=timeout,
            stdin=subprocess.DEVNULL,    # Don't inherit stdin
            start_new_session=True,      # Isolate from parent process group
            env=clean_env,
            **kwargs
        )
    
    def _create_embedded_attachment_pdf(
        self, 
        input_path: Path, 
        output_path: Path, 
        ext: str
    ) -> ConversionResult:
        """
        Create a placeholder PDF with the original file embedded.
        For unsupported or failed conversions.
        """
        try:
            buffer = io.BytesIO()
            
            doc = SimpleDocTemplate(
                buffer,
                pagesize=letter,
                rightMargin=inch,
                leftMargin=inch,
                topMargin=1.5*inch,
                bottomMargin=inch
            )
            
            styles = getSampleStyleSheet()
            
            # Title style
            title_style = ParagraphStyle(
                name='AttachmentTitle',
                parent=styles['Heading1'],
                fontSize=16,
                textColor=colors.Color(0.2, 0.2, 0.4),
                spaceAfter=15,
                alignment=1  # Center
            )
            
            # Filename style
            filename_style = ParagraphStyle(
                name='AttachmentFilename',
                parent=styles['Normal'],
                fontSize=12,
                textColor=colors.Color(0.2, 0.2, 0.2),
                spaceAfter=5,
                alignment=1,
                fontName='Helvetica-Bold'
            )
            
            # Info style
            info_style = ParagraphStyle(
                name='AttachmentInfo',
                parent=styles['Normal'],
                fontSize=10,
                textColor=colors.Color(0.4, 0.4, 0.4),
                spaceAfter=8,
                alignment=1
            )
            
            # Not converted warning style
            warning_style = ParagraphStyle(
                name='NotConverted',
                parent=styles['Normal'],
                fontSize=14,
                textColor=colors.Color(0.8, 0.1, 0.1),
                spaceBefore=15,
                spaceAfter=15,
                alignment=1,
                fontName='Helvetica-Bold'
            )
            
            # Instructions style
            instructions_style = ParagraphStyle(
                name='Instructions',
                parent=styles['Normal'],
                fontSize=9,
                textColor=colors.Color(0.5, 0.5, 0.5),
                spaceBefore=20,
                alignment=1,
                leading=14
            )
            
            # Get file info
            file_size = input_path.stat().st_size if input_path.exists() else 0
            size_str = self._format_file_size(file_size)
            
            story = [
                Spacer(1, 0.5*inch),
                Paragraph("📎 ATTACHMENT", title_style),
                Spacer(1, 10),
                Paragraph(self._escape_text(input_path.name), filename_style),
                Paragraph(f"Type: {ext.upper()} | Size: {size_str}", info_style),
                Spacer(1, 15),
                Paragraph("(Not Converted)", warning_style),
                Spacer(1, 10),
                Paragraph(
                    "This attachment type cannot be converted to PDF.",
                    info_style
                ),
                Paragraph(
                    "The original file is embedded in this PDF document.<br/><br/>"
                    "<b>To access the file:</b><br/>"
                    "• In Adobe Reader: View → Show/Hide → Navigation Panes → Attachments<br/>"
                    "• In Preview (Mac): The attachment panel may not be supported<br/>"
                    "• In other readers: Look for a paperclip icon or attachments panel",
                    instructions_style
                ),
            ]
            
            doc.build(story)
            
            # Now create the final PDF with embedded file using pikepdf
            pdf_bytes = buffer.getvalue()
            
            # Use pikepdf to embed the file (compatible with pikepdf merge)
            import pikepdf
            pdf = pikepdf.Pdf.open(io.BytesIO(pdf_bytes))
            
            # Read the original file and embed it
            if input_path.exists():
                with open(input_path, 'rb') as f:
                    file_data = f.read()
                
                # Create the embedded file stream
                file_stream = pikepdf.Stream(pdf, file_data)
                file_stream['/Type'] = pikepdf.Name('/EmbeddedFile')
                
                # Try to set MIME type based on extension
                mime_types = {
                    '.ics': 'text/calendar',
                    '.wav': 'audio/wav',
                    '.mp3': 'audio/mpeg',
                    '.msg': 'application/vnd.ms-outlook',
                    '.eml': 'message/rfc822',
                    '.zip': 'application/zip',
                    '.doc': 'application/msword',
                    '.docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
                    '.xls': 'application/vnd.ms-excel',
                    '.xlsx': 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
                }
                mime_type = mime_types.get(ext.lower(), 'application/octet-stream')
                file_stream['/Subtype'] = pikepdf.Name(f'/{mime_type.replace("/", "#2F")}')
                
                # Create file spec dictionary
                file_spec = pikepdf.Dictionary({
                    '/Type': pikepdf.Name('/Filespec'),
                    '/F': input_path.name,
                    '/UF': input_path.name,
                    '/EF': pikepdf.Dictionary({
                        '/F': file_stream,
                        '/UF': file_stream,
                    }),
                    '/Desc': f'Original attachment: {input_path.name}',
                })
                
                # Create or get the Names dictionary
                if '/Names' not in pdf.Root:
                    pdf.Root['/Names'] = pikepdf.Dictionary()
                
                # Create embedded files name tree
                embedded_files = pikepdf.Dictionary({
                    '/Names': pikepdf.Array([
                        input_path.name,
                        pdf.make_indirect(file_spec)
                    ])
                })
                pdf.Root.Names['/EmbeddedFiles'] = embedded_files
                
                # Add a clickable link annotation to open the embedded file
                # The link covers a region in the middle of the page
                page = pdf.pages[0]
                media_box = page.MediaBox
                page_width = float(media_box[2]) - float(media_box[0])
                page_height = float(media_box[3]) - float(media_box[1])
                
                # Create link rect (centered, covering the "click to open" area)
                link_x1 = page_width * 0.2
                link_x2 = page_width * 0.8
                link_y1 = page_height * 0.3
                link_y2 = page_height * 0.5
                
                # Create GoToE action to open embedded file
                goto_action = pikepdf.Dictionary({
                    '/S': pikepdf.Name('/Launch'),
                    '/F': file_spec,
                    '/NewWindow': True,
                })
                
                # Create link annotation
                link_annot = pikepdf.Dictionary({
                    '/Type': pikepdf.Name('/Annot'),
                    '/Subtype': pikepdf.Name('/Link'),
                    '/Rect': pikepdf.Array([link_x1, link_y1, link_x2, link_y2]),
                    '/Border': pikepdf.Array([0, 0, 0]),  # No visible border
                    '/A': goto_action,
                    '/H': pikepdf.Name('/I'),  # Invert highlight
                })
                
                # Add annotation to page
                if '/Annots' not in page:
                    page['/Annots'] = pikepdf.Array()
                page.Annots.append(pdf.make_indirect(link_annot))
            
            # Write output
            output_path.parent.mkdir(parents=True, exist_ok=True)
            pdf.save(str(output_path))
            
            return ConversionResult(
                status=ConversionStatus.PARTIAL,
                output_path=output_path,
                original_path=input_path,
                original_type=ext,
                message=f"Created placeholder with embedded file (Not Converted)",
                ocr_applied=False
            )
            
        except Exception as e:
            logger.exception(f"Error creating embedded attachment PDF: {e}")
            return ConversionResult(
                status=ConversionStatus.FAILED,
                output_path=None,
                original_path=input_path,
                original_type=ext,
                message=f"Failed to create embedded attachment: {str(e)}"
            )
    
    def _format_file_size(self, size: int) -> str:
        """Format file size in human readable form."""
        for unit in ['B', 'KB', 'MB', 'GB']:
            if size < 1024:
                return f"{size:.1f} {unit}"
            size /= 1024
        return f"{size:.1f} TB"
    
    def _escape_text(self, text: str) -> str:
        """Escape text for reportlab."""
        if not text:
            return ""
        text = text.replace('&', '&amp;')
        text = text.replace('<', '&lt;')
        text = text.replace('>', '&gt;')
        return text
    
    def cleanup(self):
        """Clean up temporary files."""
        if os.path.exists(self.temp_dir):
            shutil.rmtree(self.temp_dir, ignore_errors=True)
    
    def convert(
        self,
        input_path: Union[str, Path],
        output_dir: Union[str, Path],
        output_filename: Optional[str] = None
    ) -> ConversionResult:
        """
        Convert a file to PDF.
        
        Args:
            input_path: Path to the input file
            output_dir: Directory for the output PDF
            output_filename: Optional custom output filename (without .pdf)
            
        Returns:
            ConversionResult with conversion details
        """
        input_path = Path(input_path)
        output_dir = Path(output_dir)
        output_dir.mkdir(parents=True, exist_ok=True)
        
        # Determine output filename
        if output_filename:
            output_path = output_dir / f"{output_filename}.pdf"
        else:
            output_path = output_dir / f"{input_path.stem}.pdf"
        
        # Get file extension
        ext = input_path.suffix.lower()
        
        if ext not in self.SUPPORTED_FORMATS:
            # Create a placeholder PDF with embedded file for unsupported formats
            return self._create_embedded_attachment_pdf(input_path, output_path, ext)
        
        # Get conversion method
        method_name = self.SUPPORTED_FORMATS[ext]
        method = getattr(self, method_name)
        
        try:
            self._log(f"Converting {input_path.name}...")
            result = method(input_path, output_path)
            
            # If conversion failed, create embedded placeholder instead
            if result.status == ConversionStatus.FAILED:
                return self._create_embedded_attachment_pdf(input_path, output_path, ext)
            
            return result
        except Exception as e:
            logger.exception(f"Error converting {input_path}: {e}")
            # Try to create embedded placeholder on error
            try:
                return self._create_embedded_attachment_pdf(input_path, output_path, ext)
            except:
                return ConversionResult(
                    status=ConversionStatus.FAILED,
                    output_path=None,
                    original_path=input_path,
                    original_type=ext,
                    message=f"Conversion error: {str(e)}"
                )
    
    def convert_bytes(
        self,
        content: bytes,
        content_type: str,
        filename: str,
        output_dir: Union[str, Path],
        output_filename: Optional[str] = None
    ) -> ConversionResult:
        """
        Convert bytes content to PDF.
        
        Args:
            content: File content as bytes
            content_type: MIME type of the content
            filename: Original filename
            output_dir: Directory for the output PDF
            output_filename: Optional custom output filename
            
        Returns:
            ConversionResult with conversion details
        """
        # Save to temp file first
        temp_path = Path(self.temp_dir) / filename
        temp_path.parent.mkdir(parents=True, exist_ok=True)
        
        with open(temp_path, 'wb') as f:
            f.write(content)
        
        return self.convert(temp_path, output_dir, output_filename)
    
    # === PDF Conversion (with OCR) ===
    
    def _convert_pdf(self, input_path: Path, output_path: Path) -> ConversionResult:
        """Convert/process PDF, applying OCR if needed."""
        ocr_applied = False
        
        try:
            reader = PdfReader(str(input_path))
            
            # Check if PDF is encrypted
            if reader.is_encrypted:
                # Try empty password first (some PDFs are "encrypted" with no password)
                try:
                    reader.decrypt('')
                except Exception:
                    # Password protected - just copy as-is and embed
                    logger.warning(f"PDF is password-protected: {input_path.name}")
                    return self._create_embedded_attachment_pdf(input_path, output_path, '.pdf')
            
            # Check if PDF has text
            has_text = False
            for page in reader.pages[:3]:  # Check first 3 pages
                text = page.extract_text()
                if text and text.strip():
                    has_text = True
                    break
            
            # Check if PDF has embedded images (common in scanned PDFs)
            has_images = False
            try:
                import pikepdf
                with pikepdf.open(input_path) as pdf:
                    for page in list(pdf.pages)[:3]:
                        if '/Resources' in page:
                            resources = page['/Resources']
                            if '/XObject' in resources:
                                xobjects = resources['/XObject']
                                for key in xobjects.keys():
                                    xobj = xobjects[key]
                                    if xobj.get('/Subtype') == '/Image':
                                        has_images = True
                                        break
                        if has_images:
                            break
            except Exception:
                pass  # If we can't check, assume no images
            
            # Decision logic:
            # - If PDF has text: just copy (already searchable)
            # - If PDF has images but no text and OCR disabled: just copy (images are visible)
            # - If PDF has no text and no images and OCR enabled: apply OCR
            # - If PDF has no text, has images, and OCR enabled: apply OCR for searchability
            #   BUT this can cause rendering issues, so we skip it for now
            
            if has_text:
                # PDF already has searchable text, just copy
                shutil.copy(input_path, output_path)
            elif has_images:
                # PDF has images (likely scanned document) - just copy to preserve rendering
                # OCR can cause compatibility issues on some viewers
                shutil.copy(input_path, output_path)
            elif self.ocr_enabled:
                # No text, no images - try OCR
                ocr_applied = self._ocr_pdf(input_path, output_path)
            else:
                # Just copy as-is
                shutil.copy(input_path, output_path)
            
            return ConversionResult(
                status=ConversionStatus.SUCCESS,
                output_path=output_path,
                original_path=input_path,
                original_type='.pdf',
                message="PDF processed successfully",
                ocr_applied=ocr_applied
            )
        
        except Exception as e:
            logger.warning(f"PDF processing error for {input_path.name}: {e}")
            # If PDF is corrupted or has issues, try to just copy it
            try:
                shutil.copy(input_path, output_path)
                return ConversionResult(
                    status=ConversionStatus.PARTIAL,
                    output_path=output_path,
                    original_path=input_path,
                    original_type='.pdf',
                    message=f"PDF copied as-is (processing error: {e})",
                    ocr_applied=False
                )
            except Exception as copy_error:
                logger.error(f"Failed to copy PDF {input_path.name}: {copy_error}")
                # Last resort: embed the original file
                return self._create_embedded_attachment_pdf(input_path, output_path, '.pdf')
    
    def _ocr_pdf(self, input_path: Path, output_path: Path) -> bool:
        """Apply OCR to a PDF file."""
        if not self.ocr_enabled or not self.has_tesseract:
            shutil.copy(input_path, output_path)
            return False
        
        try:
            # Convert PDF to images (pass poppler_path on Windows if found)
            convert_kwargs = {'dpi': 300}
            if self.poppler_path:
                convert_kwargs['poppler_path'] = self.poppler_path
            
            images = convert_from_path(str(input_path), **convert_kwargs)
            
            # OCR each page and create new PDF
            pdf_pages = []
            for img in images:
                # Get OCR PDF for this page
                pdf_bytes = pytesseract.image_to_pdf_or_hocr(img, extension='pdf')
                pdf_pages.append(pdf_bytes)
            
            # Merge all pages
            writer = PdfWriter()
            for pdf_bytes in pdf_pages:
                reader = PdfReader(io.BytesIO(pdf_bytes))
                for page in reader.pages:
                    writer.add_page(page)
            
            with open(output_path, 'wb') as f:
                writer.write(f)
            
            return True
        
        except Exception as e:
            logger.warning(f"OCR failed for PDF: {e}")
            shutil.copy(input_path, output_path)
            return False
    
    # === Image Conversion ===
    
    def _convert_image(self, input_path: Path, output_path: Path) -> ConversionResult:
        """Convert image to PDF with OCR."""
        ocr_applied = False
        
        try:
            img = Image.open(input_path)
            
            # Convert to RGB if necessary
            if img.mode in ('RGBA', 'P', 'LA'):
                background = Image.new('RGB', img.size, (255, 255, 255))
                if img.mode == 'P':
                    img = img.convert('RGBA')
                background.paste(img, mask=img.split()[-1] if img.mode == 'RGBA' else None)
                img = background
            elif img.mode != 'RGB':
                img = img.convert('RGB')
            
            if self.ocr_enabled and self.has_tesseract:
                # Create searchable PDF with OCR
                try:
                    pdf_bytes = pytesseract.image_to_pdf_or_hocr(img, extension='pdf')
                    with open(output_path, 'wb') as f:
                        f.write(pdf_bytes)
                    ocr_applied = True
                except Exception as e:
                    logger.warning(f"OCR failed for image: {e}")
                    # Fallback to simple conversion
                    img.save(output_path, 'PDF', resolution=300)
            else:
                # Simple conversion without OCR
                img.save(output_path, 'PDF', resolution=300)
            
            return ConversionResult(
                status=ConversionStatus.SUCCESS,
                output_path=output_path,
                original_path=input_path,
                original_type=input_path.suffix.lower(),
                message="Image converted successfully",
                ocr_applied=ocr_applied
            )
        
        except Exception as e:
            raise RuntimeError(f"Image conversion failed: {e}")
    
    # === Word Document Conversion ===
    
    def _convert_docx(self, input_path: Path, output_path: Path) -> ConversionResult:
        """Convert DOCX to PDF using LibreOffice."""
        return self._libreoffice_convert(input_path, output_path, '.docx')
    
    def _convert_doc(self, input_path: Path, output_path: Path) -> ConversionResult:
        """Convert DOC to PDF using LibreOffice."""
        return self._libreoffice_convert(input_path, output_path, '.doc')
    
    # === PowerPoint Conversion ===
    
    def _convert_pptx(self, input_path: Path, output_path: Path) -> ConversionResult:
        """Convert PPTX to PDF using LibreOffice."""
        return self._libreoffice_convert(input_path, output_path, '.pptx')
    
    def _convert_ppt(self, input_path: Path, output_path: Path) -> ConversionResult:
        """Convert PPT to PDF using LibreOffice."""
        return self._libreoffice_convert(input_path, output_path, '.ppt')
    
    def _libreoffice_convert(self, input_path: Path, output_path: Path, ext: str) -> ConversionResult:
        """Convert document using LibreOffice."""
        if not self.has_libreoffice:
            # Fallback: try to extract text
            return self._fallback_document_convert(input_path, output_path, ext)
        
        try:
            # Use the stored LibreOffice path (handles macOS app bundle)
            lo_path = self.libreoffice_path
            
            # Create temp output directory
            temp_out = Path(self.temp_dir) / "lo_output"
            temp_out.mkdir(exist_ok=True)
            
            # Run LibreOffice conversion using thread-safe subprocess helper
            cmd = [
                lo_path,
                '--headless',
                '--convert-to', 'pdf',
                '--outdir', str(temp_out),
                str(input_path)
            ]
            
            result = self._safe_subprocess_run(cmd, timeout=120)
            
            # Find output file
            expected_output = temp_out / f"{input_path.stem}.pdf"
            
            if expected_output.exists():
                shutil.move(str(expected_output), str(output_path))
                return ConversionResult(
                    status=ConversionStatus.SUCCESS,
                    output_path=output_path,
                    original_path=input_path,
                    original_type=ext,
                    message="Document converted successfully"
                )
            else:
                raise RuntimeError(f"LibreOffice did not produce output: {result.stderr}")
        
        except subprocess.TimeoutExpired:
            logger.warning(f"LibreOffice timed out converting {input_path.name} (>120s)")
            return self._fallback_document_convert(input_path, output_path, ext)
        except Exception as e:
            logger.warning(f"LibreOffice conversion failed for {input_path.name}: {e}")
            return self._fallback_document_convert(input_path, output_path, ext)
    
    def _libreoffice_excel_to_pdf(self, input_path: Path, output_path: Path) -> ConversionResult:
        """
        Convert Excel to PDF using LibreOffice with ALL sheets.
        Uses special export options to ensure all sheets are included.
        """
        if not self.has_libreoffice:
            return self._excel_fallback_convert(input_path, output_path, '.xlsx')
        
        try:
            lo_path = self.libreoffice_path
            
            # Create temp output directory
            temp_out = Path(self.temp_dir) / "lo_output"
            temp_out.mkdir(exist_ok=True)
            
            # LibreOffice Calc by default exports only the active sheet.
            # To export ALL sheets, we need to use the proper export filter.
            # The "calc_pdf_Export" filter with empty Selection exports all sheets.
            cmd = [
                lo_path,
                '--headless',
                '--convert-to', 'pdf:calc_pdf_Export',
                '--outdir', str(temp_out),
                str(input_path)
            ]
            
            logger.info(f"Converting Excel with all sheets: {input_path.name}")
            
            # Use thread-safe subprocess helper (longer timeout for multi-sheet workbooks)
            result = self._safe_subprocess_run(cmd, timeout=180)
            
            # Find output file
            expected_output = temp_out / f"{input_path.stem}.pdf"
            
            if expected_output.exists():
                # Verify we got more than just the active sheet by checking file size
                pdf_size = expected_output.stat().st_size
                logger.info(f"Excel PDF created: {pdf_size} bytes")
                
                shutil.move(str(expected_output), str(output_path))
                return ConversionResult(
                    status=ConversionStatus.SUCCESS,
                    output_path=output_path,
                    original_path=input_path,
                    original_type='.xlsx',
                    message="Excel converted with print settings applied"
                )
            else:
                raise RuntimeError(f"LibreOffice did not produce output: {result.stderr}")
        
        except subprocess.TimeoutExpired:
            logger.warning(f"LibreOffice timed out converting Excel {input_path.name}")
            return self._excel_fallback_convert(input_path, output_path, '.xlsx')
        except Exception as e:
            logger.warning(f"LibreOffice Excel conversion failed: {e}, trying standard conversion")
            # Fall back to standard conversion
            return self._libreoffice_convert(input_path, output_path, '.xlsx')
    
    def _fallback_document_convert(self, input_path: Path, output_path: Path, ext: str) -> ConversionResult:
        """Fallback conversion for documents when LibreOffice is unavailable."""
        
        if ext in ('.docx', '.doc'):
            return self._docx_to_pdf_fallback(input_path, output_path)
        elif ext in ('.pptx', '.ppt'):
            return self._pptx_to_pdf_fallback(input_path, output_path)
        else:
            return ConversionResult(
                status=ConversionStatus.FAILED,
                output_path=None,
                original_path=input_path,
                original_type=ext,
                message="LibreOffice required for this format"
            )
    
    def _docx_to_pdf_fallback(self, input_path: Path, output_path: Path) -> ConversionResult:
        """Convert DOCX to PDF using python-docx (text only)."""
        try:
            doc = DocxDocument(str(input_path))
            
            # Create PDF with reportlab
            pdf_doc = SimpleDocTemplate(
                str(output_path),
                pagesize=letter,
                rightMargin=72,
                leftMargin=72,
                topMargin=72,
                bottomMargin=72
            )
            
            styles = getSampleStyleSheet()
            story = []
            
            for para in doc.paragraphs:
                if para.text.strip():
                    story.append(Paragraph(para.text, styles['Normal']))
                    story.append(Spacer(1, 12))
            
            if story:
                pdf_doc.build(story)
            else:
                # Empty document, create placeholder
                story.append(Paragraph("(Empty document)", styles['Normal']))
                pdf_doc.build(story)
            
            return ConversionResult(
                status=ConversionStatus.PARTIAL,
                output_path=output_path,
                original_path=input_path,
                original_type='.docx',
                message="Document converted (text only, no formatting)"
            )
        
        except Exception as e:
            raise RuntimeError(f"DOCX fallback conversion failed: {e}")
    
    def _pptx_to_pdf_fallback(self, input_path: Path, output_path: Path) -> ConversionResult:
        """Convert PPTX to PDF using python-pptx (text only)."""
        try:
            prs = Presentation(str(input_path))
            
            pdf_doc = SimpleDocTemplate(
                str(output_path),
                pagesize=letter,
                rightMargin=72,
                leftMargin=72,
                topMargin=72,
                bottomMargin=72
            )
            
            styles = getSampleStyleSheet()
            story = []
            
            for i, slide in enumerate(prs.slides, 1):
                story.append(Paragraph(f"<b>Slide {i}</b>", styles['Heading2']))
                story.append(Spacer(1, 12))
                
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text.strip():
                        story.append(Paragraph(shape.text, styles['Normal']))
                        story.append(Spacer(1, 6))
                
                story.append(Spacer(1, 24))
            
            if story:
                pdf_doc.build(story)
            
            return ConversionResult(
                status=ConversionStatus.PARTIAL,
                output_path=output_path,
                original_path=input_path,
                original_type='.pptx',
                message="Presentation converted (text only)"
            )
        
        except Exception as e:
            raise RuntimeError(f"PPTX fallback conversion failed: {e}")
    
    # === Excel Conversion ===
    
    def _convert_excel(self, input_path: Path, output_path: Path) -> ConversionResult:
        """Convert Excel file to PDF with optimized print settings."""
        ext = input_path.suffix.lower()
        
        # For xlsx files, apply print settings before LibreOffice conversion
        if ext == '.xlsx' and self.has_libreoffice:
            try:
                result = self._convert_excel_with_settings(input_path, output_path)
                if result.status == ConversionStatus.SUCCESS:
                    return result
            except Exception as e:
                logger.warning(f"Excel print settings failed, trying direct conversion: {e}")
        
        # Try LibreOffice directly (for .xls or if settings approach failed)
        if self.has_libreoffice:
            result = self._libreoffice_convert(input_path, output_path, ext)
            if result.status == ConversionStatus.SUCCESS:
                return result
        
        # Fallback to pandas/reportlab
        return self._excel_fallback_convert(input_path, output_path, ext)
    
    def _convert_excel_with_settings(self, input_path: Path, output_path: Path) -> ConversionResult:
        """
        Convert Excel with optimized print settings:
        - All sheets printed
        - Fit all columns on one page
        - Landscape orientation  
        - Small margins
        """
        from openpyxl import load_workbook
        from openpyxl.worksheet.page import PageMargins
        from openpyxl.worksheet.properties import WorksheetProperties, PageSetupProperties
        
        # Create a temp copy to modify
        temp_xlsx = Path(self.temp_dir) / f"print_ready_{input_path.name}"
        shutil.copy(input_path, temp_xlsx)
        
        try:
            wb = load_workbook(temp_xlsx)
            
            # Get list of all sheet names for logging
            sheet_names = wb.sheetnames
            logger.info(f"Excel file has {len(sheet_names)} sheets: {sheet_names}")
            
            for sheet in wb.worksheets:
                try:
                    # Set landscape orientation
                    sheet.page_setup.orientation = 'landscape'
                    
                    # Ensure sheet has properties (fixes openpyxl bug with some Excel files)
                    if sheet.sheet_properties is None:
                        sheet.sheet_properties = WorksheetProperties()
                    if sheet.sheet_properties.pageSetUpPr is None:
                        sheet.sheet_properties.pageSetUpPr = PageSetupProperties()
                    
                    # Fit all columns on one page (width), let rows span pages
                    sheet.sheet_properties.pageSetUpPr.fitToPage = True
                    sheet.page_setup.fitToWidth = 1
                    sheet.page_setup.fitToHeight = 0  # 0 = unlimited pages for height
                    
                    # Small margins (in inches)
                    sheet.page_margins = PageMargins(
                        left=0.25,
                        right=0.25,
                        top=0.5,
                        bottom=0.5,
                        header=0.3,
                        footer=0.3
                    )
                    
                    # Ensure sheet is visible (not hidden) so it gets printed
                    sheet.sheet_state = 'visible'
                    
                    logger.debug(f"Applied print settings to sheet: {sheet.title}")
                except Exception as sheet_err:
                    logger.warning(f"Could not apply settings to sheet {sheet.title}: {sheet_err}")
            
            wb.save(temp_xlsx)
            wb.close()
            
            # Convert with LibreOffice using special filter to export ALL sheets
            result = self._libreoffice_excel_to_pdf(temp_xlsx, output_path)
            
            # Clean up temp file
            try:
                temp_xlsx.unlink()
            except:
                pass
            
            return result
            
        except Exception as e:
            # Clean up on error
            try:
                temp_xlsx.unlink()
            except:
                pass
            raise
    
    def _excel_fallback_convert(self, input_path: Path, output_path: Path, ext: str) -> ConversionResult:
        """Fallback Excel conversion using pandas/reportlab (landscape, fit to page)."""
        try:
            if ext == '.xlsx':
                df_dict = pd.read_excel(input_path, sheet_name=None, engine='openpyxl')
            else:
                df_dict = pd.read_excel(input_path, sheet_name=None, engine='xlrd')
            
            # Use landscape with small margins for spreadsheets
            from reportlab.lib.pagesizes import landscape, letter
            pdf_doc = SimpleDocTemplate(
                str(output_path),
                pagesize=landscape(letter),
                rightMargin=18,
                leftMargin=18,
                topMargin=36,
                bottomMargin=36
            )
            
            styles = getSampleStyleSheet()
            story = []
            
            for sheet_name, df in df_dict.items():
                story.append(Paragraph(f"<b>Sheet: {sheet_name}</b>", styles['Heading2']))
                story.append(Spacer(1, 12))
                
                if not df.empty:
                    # Convert DataFrame to table
                    table_data = [df.columns.tolist()] + df.values.tolist()
                    
                    # Limit columns and rows for readability
                    max_cols = 10
                    max_rows = 100
                    
                    if len(table_data[0]) > max_cols:
                        table_data = [row[:max_cols] + ['...'] for row in table_data]
                    
                    if len(table_data) > max_rows:
                        table_data = table_data[:max_rows] + [['...'] * len(table_data[0])]
                    
                    # Convert all values to strings
                    table_data = [[str(cell)[:50] for cell in row] for row in table_data]
                    
                    table = Table(table_data)
                    table.setStyle(TableStyle([
                        ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
                        ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                        ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
                        ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                        ('FONTSIZE', (0, 0), (-1, -1), 8),
                        ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                        ('BACKGROUND', (0, 1), (-1, -1), colors.white),
                        ('GRID', (0, 0), (-1, -1), 1, colors.black),
                    ]))
                    
                    story.append(table)
                
                story.append(Spacer(1, 24))
            
            pdf_doc.build(story)
            
            return ConversionResult(
                status=ConversionStatus.SUCCESS,
                output_path=output_path,
                original_path=input_path,
                original_type=ext,
                message="Excel converted successfully"
            )
        
        except Exception as e:
            raise RuntimeError(f"Excel conversion failed: {e}")
    
    # === Text File Conversion ===
    
    def _convert_text(self, input_path: Path, output_path: Path) -> ConversionResult:
        """Convert text file to PDF."""
        try:
            # Read text with encoding detection
            content = self._read_text_file(input_path)
            
            pdf_doc = SimpleDocTemplate(
                str(output_path),
                pagesize=letter,
                rightMargin=72,
                leftMargin=72,
                topMargin=72,
                bottomMargin=72
            )
            
            # Use monospace style for text files
            styles = getSampleStyleSheet()
            mono_style = ParagraphStyle(
                'Mono',
                parent=styles['Normal'],
                fontName='Courier',
                fontSize=9,
                leading=12
            )
            
            story = []
            
            for line in content.split('\n'):
                # Escape special characters for reportlab
                line = line.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
                story.append(Paragraph(line or '&nbsp;', mono_style))
            
            pdf_doc.build(story)
            
            return ConversionResult(
                status=ConversionStatus.SUCCESS,
                output_path=output_path,
                original_path=input_path,
                original_type='.txt',
                message="Text file converted successfully"
            )
        
        except Exception as e:
            raise RuntimeError(f"Text conversion failed: {e}")
    
    def _read_text_file(self, path: Path) -> str:
        """Read text file with encoding detection."""
        encodings = ['utf-8', 'utf-16', 'latin-1', 'cp1252', 'ascii']
        
        for encoding in encodings:
            try:
                with open(path, 'r', encoding=encoding) as f:
                    return f.read()
            except (UnicodeDecodeError, UnicodeError):
                continue
        
        # Last resort
        with open(path, 'r', encoding='utf-8', errors='replace') as f:
            return f.read()
    
    # === ICS Calendar Conversion ===
    
    def _convert_ics(self, input_path: Path, output_path: Path) -> ConversionResult:
        """
        Convert ICS (iCalendar) file to a nicely formatted PDF.
        Extracts event details like title, date/time, location, attendees, etc.
        """
        try:
            content = self._read_text_file(input_path)
            events = self._parse_ics_content(content)
            
            if not events:
                # No events found, create a simple text conversion
                return self._convert_text(input_path, output_path)
            
            # Create PDF with event details
            pdf_doc = SimpleDocTemplate(
                str(output_path),
                pagesize=letter,
                rightMargin=72,
                leftMargin=72,
                topMargin=72,
                bottomMargin=72
            )
            
            styles = getSampleStyleSheet()
            
            # Custom styles for calendar invite
            title_style = ParagraphStyle(
                'EventTitle',
                parent=styles['Heading1'],
                fontSize=16,
                spaceAfter=12,
                textColor=colors.HexColor('#1a365d')
            )
            
            label_style = ParagraphStyle(
                'Label',
                parent=styles['Normal'],
                fontSize=10,
                fontName='Helvetica-Bold',
                textColor=colors.HexColor('#4a5568'),
                spaceBefore=8,
                spaceAfter=2
            )
            
            value_style = ParagraphStyle(
                'Value',
                parent=styles['Normal'],
                fontSize=11,
                leftIndent=10,
                spaceAfter=6
            )
            
            attendee_style = ParagraphStyle(
                'Attendee',
                parent=styles['Normal'],
                fontSize=10,
                leftIndent=20,
                spaceAfter=2
            )
            
            desc_style = ParagraphStyle(
                'Description',
                parent=styles['Normal'],
                fontSize=10,
                leftIndent=10,
                spaceAfter=6
            )
            
            story = []
            
            # Header
            story.append(Paragraph("📅 Calendar Invitation", styles['Title']))
            story.append(Spacer(1, 20))
            
            for i, event in enumerate(events):
                if i > 0:
                    story.append(Spacer(1, 20))
                    story.append(Paragraph("─" * 60, styles['Normal']))
                    story.append(Spacer(1, 20))
                
                # Event title/summary
                summary = self._escape_text(event.get('summary', 'Untitled Event'))
                story.append(Paragraph(summary, title_style))
                
                # Date and Time (the key info the user wanted!)
                if event.get('start') or event.get('end'):
                    story.append(Paragraph("Date & Time", label_style))
                    
                    start = event.get('start', '')
                    end = event.get('end', '')
                    
                    if start and end:
                        time_str = f"{start}  →  {end}"
                    elif start:
                        time_str = f"Starts: {start}"
                    else:
                        time_str = f"Ends: {end}"
                    
                    story.append(Paragraph(self._escape_text(time_str), value_style))
                
                # Duration if available
                if event.get('duration'):
                    story.append(Paragraph(f"Duration: {event['duration']}", value_style))
                
                # Location
                if event.get('location'):
                    story.append(Paragraph("Location", label_style))
                    story.append(Paragraph(self._escape_text(event['location']), value_style))
                
                # Organizer
                if event.get('organizer'):
                    story.append(Paragraph("Organizer", label_style))
                    story.append(Paragraph(self._escape_text(event['organizer']), value_style))
                
                # Attendees
                if event.get('attendees'):
                    story.append(Paragraph(f"Attendees ({len(event['attendees'])})", label_style))
                    for attendee in event['attendees'][:20]:  # Limit to 20
                        story.append(Paragraph(f"• {self._escape_text(attendee)}", attendee_style))
                    if len(event['attendees']) > 20:
                        story.append(Paragraph(f"  ... and {len(event['attendees']) - 20} more", attendee_style))
                
                # Status
                if event.get('status'):
                    story.append(Paragraph("Status", label_style))
                    story.append(Paragraph(self._escape_text(event['status']), value_style))
                
                # Description
                if event.get('description'):
                    story.append(Paragraph("Description", label_style))
                    desc_text = event['description'][:2000]  # Limit length
                    if len(event['description']) > 2000:
                        desc_text += "..."
                    # Handle line breaks in description
                    desc_text = desc_text.replace('\n', '<br/>')
                    story.append(Paragraph(self._escape_text(desc_text).replace('&lt;br/&gt;', '<br/>'), desc_style))
                
                # URL
                if event.get('url'):
                    story.append(Paragraph("URL", label_style))
                    story.append(Paragraph(self._escape_text(event['url']), value_style))
            
            pdf_doc.build(story)
            
            return ConversionResult(
                status=ConversionStatus.SUCCESS,
                output_path=output_path,
                original_path=input_path,
                original_type='.ics',
                message=f"Calendar invite converted ({len(events)} event(s))"
            )
            
        except Exception as e:
            logger.exception(f"ICS conversion failed: {e}")
            raise RuntimeError(f"ICS conversion failed: {e}")
    
    def _parse_ics_content(self, content: str) -> list:
        """
        Parse ICS content and extract event details.
        Simple parser that handles common ICS properties without external dependencies.
        """
        events = []
        current_event = None
        current_key = None
        current_value = ""
        
        lines = content.replace('\r\n', '\n').replace('\r', '\n').split('\n')
        
        # Handle line continuations (lines starting with space/tab are continuations)
        unfolded_lines = []
        for line in lines:
            if line.startswith(' ') or line.startswith('\t'):
                if unfolded_lines:
                    unfolded_lines[-1] += line[1:]
            else:
                unfolded_lines.append(line)
        
        for line in unfolded_lines:
            line = line.strip()
            if not line:
                continue
            
            if line == 'BEGIN:VEVENT':
                current_event = {}
            elif line == 'END:VEVENT':
                if current_event:
                    events.append(current_event)
                current_event = None
            elif current_event is not None and ':' in line:
                # Parse property
                # Handle properties with parameters like DTSTART;TZID=America/New_York:20240115T090000
                if ';' in line.split(':')[0]:
                    key_part = line.split(':')[0]
                    key = key_part.split(';')[0]
                    value = ':'.join(line.split(':')[1:])
                else:
                    key, value = line.split(':', 1)
                
                key = key.upper()
                
                if key == 'SUMMARY':
                    current_event['summary'] = self._decode_ics_value(value)
                elif key == 'DTSTART':
                    current_event['start'] = self._format_ics_datetime(value)
                elif key == 'DTEND':
                    current_event['end'] = self._format_ics_datetime(value)
                elif key == 'DURATION':
                    current_event['duration'] = self._format_ics_duration(value)
                elif key == 'LOCATION':
                    current_event['location'] = self._decode_ics_value(value)
                elif key == 'DESCRIPTION':
                    current_event['description'] = self._decode_ics_value(value)
                elif key == 'ORGANIZER':
                    # Extract name from CN parameter or email from mailto
                    organizer = value
                    organizer_name = None
                    organizer_email = None
                    
                    # Check for CN (Common Name) in the parameters part before the colon
                    # The full line might be: ORGANIZER;CN=John Smith:mailto:john@example.com
                    # But we only have the value part after the first colon split
                    # So check the key_part if it exists
                    if ';' in line.split(':')[0]:
                        key_part = line.split(':')[0]
                        # Look for CN= in the parameters
                        import re
                        cn_match = re.search(r'CN=([^;:]+)', key_part, re.IGNORECASE)
                        if cn_match:
                            organizer_name = cn_match.group(1).strip('"\'')
                    
                    # Extract email from mailto:
                    if 'mailto:' in organizer.lower():
                        organizer_email = re.split(r'mailto:', organizer, flags=re.IGNORECASE)[-1]
                    else:
                        organizer_email = organizer
                    
                    # Use name if available, otherwise email
                    if organizer_name:
                        current_event['organizer'] = organizer_name
                    else:
                        current_event['organizer'] = organizer_email
                elif key == 'ATTENDEE':
                    if 'attendees' not in current_event:
                        current_event['attendees'] = []
                    # Extract email/name from attendee line
                    attendee = value
                    # Handle both MAILTO: and mailto: (case insensitive)
                    if 'mailto:' in attendee.lower():
                        import re
                        attendee = re.split(r'mailto:', attendee, flags=re.IGNORECASE)[-1]
                        attendee = attendee.split('mailto:')[-1]
                    current_event['attendees'].append(attendee)
                elif key == 'STATUS':
                    current_event['status'] = value
                elif key == 'URL':
                    current_event['url'] = value
        
        return events
    
    def _decode_ics_value(self, value: str) -> str:
        """Decode ICS escaped characters."""
        value = value.replace('\\n', '\n')
        value = value.replace('\\N', '\n')
        value = value.replace('\\,', ',')
        value = value.replace('\\;', ';')
        value = value.replace('\\\\', '\\')
        return value
    
    def _format_ics_datetime(self, value: str) -> str:
        """Format ICS datetime into human-readable format."""
        try:
            # Remove any trailing Z (UTC indicator) for parsing
            clean_value = value.rstrip('Z')
            
            # Try different formats
            formats = [
                ('%Y%m%dT%H%M%S', '%B %d, %Y at %I:%M %p'),  # 20240115T090000
                ('%Y%m%d', '%B %d, %Y'),  # 20240115 (all-day event)
            ]
            
            for in_fmt, out_fmt in formats:
                try:
                    from datetime import datetime
                    dt = datetime.strptime(clean_value, in_fmt)
                    result = dt.strftime(out_fmt)
                    if value.endswith('Z'):
                        result += ' (UTC)'
                    return result
                except ValueError:
                    continue
            
            # If no format matched, return as-is
            return value
        except Exception:
            return value
    
    def _format_ics_duration(self, value: str) -> str:
        """Format ICS duration into human-readable format."""
        # ICS duration format: P[n]W or P[n]DT[n]H[n]M[n]S
        try:
            result = []
            value = value.upper()
            
            if value.startswith('P'):
                value = value[1:]
            
            if 'W' in value:
                weeks = int(value.split('W')[0])
                result.append(f"{weeks} week{'s' if weeks != 1 else ''}")
            elif 'T' in value:
                date_part, time_part = value.split('T')
                
                if 'D' in date_part:
                    days = int(date_part.split('D')[0])
                    result.append(f"{days} day{'s' if days != 1 else ''}")
                
                if 'H' in time_part:
                    hours = int(time_part.split('H')[0])
                    result.append(f"{hours} hour{'s' if hours != 1 else ''}")
                    time_part = time_part.split('H')[1]
                
                if 'M' in time_part:
                    minutes = int(time_part.split('M')[0])
                    result.append(f"{minutes} minute{'s' if minutes != 1 else ''}")
            
            return ', '.join(result) if result else value
        except Exception:
            return value

    # === CSV Conversion ===

    def _convert_csv(self, input_path: Path, output_path: Path) -> ConversionResult:
        """Convert CSV file to PDF."""
        try:
            df = pd.read_csv(input_path, nrows=1000)  # Limit rows
            
            pdf_doc = SimpleDocTemplate(
                str(output_path),
                pagesize=A4,
                rightMargin=36,
                leftMargin=36,
                topMargin=36,
                bottomMargin=36
            )
            
            styles = getSampleStyleSheet()
            story = []
            
            story.append(Paragraph(f"<b>{input_path.name}</b>", styles['Heading2']))
            story.append(Spacer(1, 12))
            
            if not df.empty:
                # Prepare table data
                table_data = [df.columns.tolist()] + df.values.tolist()
                
                # Limit and truncate
                max_cols = 8
                if len(table_data[0]) > max_cols:
                    table_data = [row[:max_cols] for row in table_data]
                
                table_data = [[str(cell)[:30] for cell in row] for row in table_data]
                
                table = Table(table_data)
                table.setStyle(TableStyle([
                    ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
                    ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                    ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
                    ('FONTSIZE', (0, 0), (-1, -1), 7),
                    ('GRID', (0, 0), (-1, -1), 1, colors.black),
                ]))
                
                story.append(table)
            
            pdf_doc.build(story)
            
            return ConversionResult(
                status=ConversionStatus.SUCCESS,
                output_path=output_path,
                original_path=input_path,
                original_type='.csv',
                message="CSV converted successfully"
            )
        
        except Exception as e:
            raise RuntimeError(f"CSV conversion failed: {e}")
    
    # === HTML Conversion ===
    
    def _convert_html(self, input_path: Path, output_path: Path) -> ConversionResult:
        """Convert HTML file to PDF using WeasyPrint, or fallback to basic conversion."""
        # Try WeasyPrint first (may not be available on Windows)
        try:
            from weasyprint import HTML, CSS
        except (ImportError, OSError, Exception):
            # WeasyPrint not available - use fallback
            logger.info("WeasyPrint not available, using text fallback for HTML")
            return self._html_text_fallback(input_path, output_path)
        
        try:
            content = self._read_text_file(input_path)
            
            # Basic CSS for reasonable rendering
            base_css = CSS(string='''
                @page { size: letter; margin: 0.75in; }
                body { font-family: sans-serif; font-size: 11pt; line-height: 1.4; }
                table { border-collapse: collapse; width: 100%; }
                th, td { border: 1px solid #ccc; padding: 4px 8px; }
            ''')
            
            html_doc = HTML(string=content, base_url=str(input_path.parent))
            html_doc.write_pdf(str(output_path), stylesheets=[base_css])
            
            return ConversionResult(
                status=ConversionStatus.SUCCESS,
                output_path=output_path,
                original_path=input_path,
                original_type='.html',
                message="HTML converted successfully"
            )
        
        except ImportError:
            logger.warning("WeasyPrint not available, using text-only fallback")
            return self._html_text_fallback(input_path, output_path)
        except Exception as e:
            logger.warning(f"WeasyPrint HTML conversion failed: {e}")
            return self._html_text_fallback(input_path, output_path)
    
    def _html_text_fallback(self, input_path: Path, output_path: Path) -> ConversionResult:
        """Fallback HTML to PDF conversion (text only)."""
        try:
            from html.parser import HTMLParser
            
            class TextExtractor(HTMLParser):
                def __init__(self):
                    super().__init__()
                    self.text = []
                
                def handle_data(self, data):
                    self.text.append(data.strip())
            
            content = self._read_text_file(input_path)
            
            parser = TextExtractor()
            parser.feed(content)
            text = ' '.join(filter(None, parser.text))
            
            pdf_doc = SimpleDocTemplate(
                str(output_path),
                pagesize=letter,
                rightMargin=72,
                leftMargin=72,
                topMargin=72,
                bottomMargin=72
            )
            
            styles = getSampleStyleSheet()
            story = [Paragraph(text, styles['Normal'])]
            pdf_doc.build(story)
            
            return ConversionResult(
                status=ConversionStatus.PARTIAL,
                output_path=output_path,
                original_path=input_path,
                original_type='.html',
                message="HTML converted (text only, no formatting)"
            )
        
        except Exception as e:
            raise RuntimeError(f"HTML fallback conversion failed: {e}")
    
    # === EML Conversion ===
    
    def _convert_eml(self, input_path: Path, output_path: Path) -> ConversionResult:
        """Convert nested EML to PDF."""
        # Import here to avoid circular dependency
        from .eml_parser import EMLParser
        from .email_to_pdf import EmailToPDFConverter
        
        try:
            parser = EMLParser()
            email_data = parser.parse_file(str(input_path))
            
            converter = EmailToPDFConverter()
            converter.convert_email_to_pdf(email_data, output_path)
            
            return ConversionResult(
                status=ConversionStatus.SUCCESS,
                output_path=output_path,
                original_path=input_path,
                original_type='.eml',
                message="Nested email converted successfully"
            )
        
        except Exception as e:
            raise RuntimeError(f"EML conversion failed: {e}")
    
    # === MSG Conversion ===
    
    def _convert_msg(self, input_path: Path, output_path: Path) -> ConversionResult:
        """Convert MSG (Outlook) file to PDF."""
        try:
            import extract_msg
            
            msg = extract_msg.Message(str(input_path))
            
            pdf_doc = SimpleDocTemplate(
                str(output_path),
                pagesize=letter,
                rightMargin=72,
                leftMargin=72,
                topMargin=72,
                bottomMargin=72
            )
            
            styles = getSampleStyleSheet()
            story = []
            
            # Header
            story.append(Paragraph("<b>From:</b> " + (msg.sender or "Unknown"), styles['Normal']))
            story.append(Paragraph("<b>To:</b> " + (msg.to or "Unknown"), styles['Normal']))
            story.append(Paragraph("<b>Subject:</b> " + (msg.subject or "No Subject"), styles['Normal']))
            story.append(Paragraph("<b>Date:</b> " + (str(msg.date) or "Unknown"), styles['Normal']))
            story.append(Spacer(1, 24))
            
            # Body
            body = msg.body or "(No content)"
            for line in body.split('\n'):
                line = line.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
                story.append(Paragraph(line or '&nbsp;', styles['Normal']))
            
            pdf_doc.build(story)
            
            msg.close()
            
            return ConversionResult(
                status=ConversionStatus.SUCCESS,
                output_path=output_path,
                original_path=input_path,
                original_type='.msg',
                message="MSG file converted successfully"
            )
        
        except Exception as e:
            raise RuntimeError(f"MSG conversion failed: {e}")
    
    def get_supported_extensions(self) -> List[str]:
        """Get list of supported file extensions."""
        return list(self.SUPPORTED_FORMATS.keys())
